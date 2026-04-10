import os
from django.conf import settings
import fitz  # PyMuPDF
from pptx import Presentation
from docx import Document

class DocumentService:
    @staticmethod
    def extract_text(file_obj, filename):
        """
        Extrae texto de archivos PDF, PPTX y DOCX alojados en memoria o disco temporal.
        Retorna un string con todo el texto extraído.
        """
        temp_path = os.path.join(settings.FILE_UPLOAD_TEMP_DIR, filename)
        
        # Guardar en disco temporal para procesar uniformemente
        with open(temp_path, 'wb+') as dest:
            for chunk in file_obj.chunks():
                dest.write(chunk)
                
        text = ""
        ext = filename.lower().split('.')[-1]

        try:
            if ext == 'pdf':
                text = DocumentService._extract_pdf(temp_path)
            elif ext == 'pptx':
                text = DocumentService._extract_pptx(temp_path)
            elif ext == 'docx':
                text = DocumentService._extract_docx(temp_path)
            else:
                raise ValueError(f"Formato no soportado: {ext}")
        finally:
            # Limpiar archivo temporal
            if os.path.exists(temp_path):
                os.remove(temp_path)

        return text

    @staticmethod
    def _extract_pdf(filepath):
        text = ""
        doc = fitz.open(filepath)
        for page in doc:
            text += page.get_text() + "\n"
        return text

    @staticmethod
    def _extract_pptx(filepath):
        text = ""
        prs = Presentation(filepath)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text

    @staticmethod
    def _extract_docx(filepath):
        doc = Document(filepath)
        return "\n".join([para.text for para in doc.paragraphs])
